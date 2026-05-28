"""Build an editable PowerPoint version of pitch-deck.html.

Each slide mirrors a section in pitch-deck.html. Everything is real PPTX
shapes (text boxes, tables, rectangles) so the user can edit type, colour,
and layout natively in PowerPoint or Keynote.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree

# --- Brand tokens ---------------------------------------------------------
INK      = RGBColor(0x1D, 0x1D, 0x1F)
INK2     = RGBColor(0x6E, 0x6E, 0x73)
INK3     = RGBColor(0x86, 0x86, 0x8B)
PAPER    = RGBColor(0xFB, 0xFB, 0xFD)
SURFACE  = RGBColor(0xFF, 0xFF, 0xFF)
SURFACE2 = RGBColor(0xF5, 0xF5, 0xF7)
RULE     = RGBColor(0xD2, 0xD2, 0xD7)
BLUE     = RGBColor(0x00, 0x71, 0xE3)

FONT = "Inter"  # Falls back to Calibri in PowerPoint if Inter is not installed

# --- Geometry -------------------------------------------------------------
# Slide is 16:9 wide-deck — 13.333 x 7.5 inches (1280 x 720 px @ 96dpi).
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
# HTML uses 1440x900 with 40/128/32 padding. Scale to 13.333" / 7.5".
MARGIN_X = Inches(1.2)
MARGIN_TOP = Inches(0.35)
MARGIN_BOTTOM = Inches(0.4)
CONTENT_W = SLIDE_W - 2 * MARGIN_X


def make_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank_layout)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = PAPER
    bg.shadow.inherit = False
    # remove default shadow
    spPr = bg.fill._xPr
    return s


def add_text(slide, left, top, width, height,
             text, *, font=FONT, size=14, bold=False, italic=False,
             color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.2, letter_spacing=None, uppercase=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor

    lines = text if isinstance(text, list) else [text]
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing

        runs = line if isinstance(line, list) else [{"t": line}]
        for ridx, r in enumerate(runs):
            run = p.add_run() if (idx == 0 and ridx == 0 and p.runs) else p.add_run()
            t = r.get("t", "")
            if uppercase:
                t = t.upper()
            run.text = t
            run.font.name = r.get("font", font)
            run.font.size = Pt(r.get("size", size))
            run.font.bold = r.get("bold", bold)
            run.font.italic = r.get("italic", italic)
            c = r.get("color", color)
            run.font.color.rgb = c
            ls = r.get("letter_spacing", letter_spacing)
            if ls is not None:
                rPr = run._r.get_or_add_rPr()
                rPr.set("spc", str(int(ls)))
    return tb


def add_rect(slide, left, top, width, height,
             fill=SURFACE, line=None, line_w=None, radius=None):
    if radius:
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        # The adjustment value: radius proportion (0..0.5 of shorter side)
        shorter = min(width, height)
        adj = max(0.01, min(0.5, radius / shorter))
        shp.adjustments[0] = adj
    else:
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        if line_w is not None:
            shp.line.width = line_w
    shp.shadow.inherit = False
    return shp


def add_hline(slide, left, top, width, color=RULE, weight=Pt(0.75)):
    ln = slide.shapes.add_connector(1, left, top, left + width, top)
    ln.line.color.rgb = color
    ln.line.width = weight
    return ln


# --- Common slide chrome --------------------------------------------------

def add_chrome(slide, act_text):
    add_text(slide, MARGIN_X, MARGIN_TOP, Inches(4), Inches(0.3),
             "VELO WORKING",
             size=10, bold=True, color=INK, letter_spacing=160)
    add_text(slide, SLIDE_W - MARGIN_X - Inches(4.5), MARGIN_TOP,
             Inches(4.5), Inches(0.3),
             act_text.upper(),
             size=10, bold=True, color=INK2, letter_spacing=160, align=PP_ALIGN.RIGHT)


def add_footer(slide, tagline, page_num, total=20, hl=True):
    y = SLIDE_H - MARGIN_BOTTOM
    add_text(slide, MARGIN_X, y, Inches(8), Inches(0.3),
             tagline.upper(),
             size=9, bold=True, color=INK3, letter_spacing=160)
    page_text = [{"t": f"{page_num:02d} ", "color": BLUE if hl else INK, "bold": True, "size": 10},
                 {"t": "/ 20", "color": INK3, "size": 10, "letter_spacing": 160, "bold": True}]
    add_text(slide, SLIDE_W - MARGIN_X - Inches(2), y, Inches(2), Inches(0.3),
             [page_text], align=PP_ALIGN.RIGHT)


def add_eyebrow(slide, top, text):
    # Short blue accent rule above eyebrow
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 MARGIN_X, top, Inches(0.6), Pt(1.5))
    bar.fill.solid(); bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()
    add_text(slide, MARGIN_X, top + Inches(0.18), Inches(10), Inches(0.3),
             text.upper(), size=11, bold=True, color=BLUE, letter_spacing=160)


# --- Slide builders -------------------------------------------------------

def slide_00_cover(prs):
    s = blank_slide(prs)
    add_chrome(s, "Pitch Deck · v3 · April 2026")
    # Subhead
    add_text(s, MARGIN_X, Inches(1.1), Inches(11), Inches(0.4),
             "FOR SMES WITH NON-CONTINUOUS AND DISCRETE BUSINESS SYSTEMS",
             size=10, bold=True, color=INK2, letter_spacing=200)
    # Headline
    add_text(s, MARGIN_X, Inches(1.65), Inches(11), Inches(3.0),
             [[{"t": "Own the system", "size": 64, "bold": True, "color": INK}],
              [{"t": "that runs", "size": 64, "bold": True, "color": INK}],
              [{"t": "your business.", "size": 64, "bold": True, "italic": True, "color": BLUE}]],
             line_spacing=1.05)
    # Description
    add_text(s, MARGIN_X, Inches(4.85), Inches(10), Inches(1.0),
             [[{"t": "No lock-in. No subscription. No vendor dependency."}],
              [{"t": "Cost-efficient software and hardware, self-manageable by your own employees."}]],
             size=14, color=INK2, line_spacing=1.4)
    add_footer(s, "Built for the ones still standing.", 0, hl=False)


def slide_01_about(prs):
    s = blank_slide(prs)
    add_chrome(s, "About · 01")
    add_eyebrow(s, Inches(0.85), "About Us")
    # Main paragraph
    add_text(s, MARGIN_X, Inches(1.5), Inches(11), Inches(2.6),
             [
                 [
                     {"t": "Velo Working helps SMEs with non-continuous and discrete business systems", "bold": True},
                     {"t": " — providing a full solution in automating and digitising their operations, by setting up "},
                     {"t": "no-code-based software and hardware", "bold": True},
                     {"t": " that is "},
                     {"t": "cost-efficient", "bold": True},
                     {"t": " and "},
                     {"t": "self-manageable by their own employees", "bold": True},
                     {"t": "."},
                 ]
             ], size=22, color=INK, line_spacing=1.3)

    # Left: result quotes
    add_text(s, MARGIN_X, Inches(4.7), Inches(5.5), Inches(0.4),
             "THE RESULT — IN THE OWNER'S OWN WORDS",
             size=10, bold=True, color=BLUE, letter_spacing=160)
    add_text(s, MARGIN_X, Inches(5.05), Inches(5.5), Inches(1.7),
             [[{"t": "Stops finding out problems after the damage is done.", "italic": True}],
              [{"t": "Stops guessing if the numbers are real.", "italic": True}],
              [{"t": "Stops losing everything every time someone walks out the door.", "italic": True}]],
             size=14, color=INK2, line_spacing=1.4)

    # Right: blue quote box
    bx_l = MARGIN_X + Inches(6.2); bx_t = Inches(4.7); bx_w = Inches(4.7); bx_h = Inches(1.8)
    add_rect(s, bx_l, bx_t, bx_w, bx_h, fill=SURFACE2, radius=Emu(228600))
    add_rect(s, bx_l, bx_t, Pt(3), bx_h, fill=BLUE)
    add_text(s, bx_l + Inches(0.3), bx_t + Inches(0.25), bx_w - Inches(0.5), bx_h - Inches(0.4),
             "We sell the labour. You own everything else — the software, the data, the server, the workflows.",
             size=14, italic=True, color=BLUE, line_spacing=1.3)

    add_footer(s, "One sentence. The whole positioning.", 1)


def slide_02_economy(prs):
    s = blank_slide(prs)
    add_chrome(s, "Act 1 · Ordinary World · 02")
    add_eyebrow(s, Inches(0.85), "The ground beneath the economy")
    add_text(s, MARGIN_X, Inches(1.5), Inches(11), Inches(2.2),
             [[{"t": "SMEs aren't a footnote", "size": 42, "bold": True, "color": INK}],
              [{"t": "to the economy.", "size": 42, "bold": True, "color": INK}],
              [{"t": "They are the economy.", "size": 42, "bold": True, "italic": True, "color": BLUE}]],
             line_spacing=1.05)

    stats = [
        ("99%", "of all enterprises in Singapore are SMEs.", "Enterprise Singapore"),
        ("70%", "of Singapore's workforce is employed by SMEs.", "MTI 2025"),
        ("97%", "of all enterprises in Malaysia are SMEs.", "SME Corp Malaysia"),
        ("~48%", "combined GDP contribution across Singapore and Malaysia.", "MTI / BNM 2025"),
    ]
    col_w = CONTENT_W / 4
    y0 = Inches(4.7)
    for i, (pct, label, src) in enumerate(stats):
        x = MARGIN_X + col_w * i
        add_text(s, x, y0, col_w - Inches(0.2), Inches(1.0), pct,
                 size=64, bold=True, color=INK)
        add_text(s, x, y0 + Inches(1.0), col_w - Inches(0.2), Inches(0.9), label,
                 size=11, color=INK2, line_spacing=1.4)
        add_text(s, x, y0 + Inches(1.9), col_w - Inches(0.2), Inches(0.3), src.upper(),
                 size=9, bold=True, color=INK3, letter_spacing=160)
    add_footer(s, "Precision machinists. Second-gen traders. Stamping shops. The engine of the economy.", 2)


def slide_03_paradox(prs):
    s = blank_slide(prs)
    add_chrome(s, "Act 1 · The Paradox · 03")
    add_eyebrow(s, Inches(0.85), "Research · Why digitalisation rarely sticks")
    add_text(s, MARGIN_X, Inches(1.5), Inches(11), Inches(1.8),
             [[{"t": "They want it.", "size": 40, "bold": True, "color": INK}],
              [{"t": "Very few make it stick.", "size": 40, "bold": True, "italic": True, "color": BLUE}]],
             line_spacing=1.05)
    stats = [
        ("83%", "of SMEs in Singapore have a digital transformation strategy in place.", "ASME · SME Digital Study"),
        ("~40%", "of those SMEs consider their digital efforts a success.", "ASME · SME Digital Study"),
        ("70–85%", "of ERP projects fail on cost, schedule, or fit with the actual business.", "Peer-reviewed ERP research"),
        ("58%", "of SG businesses cite high cost of technology adoption as the primary barrier.", "SBF Business Survey 2025"),
    ]
    col_w = CONTENT_W / 4
    y0 = Inches(3.85)
    for i, (pct, label, src) in enumerate(stats):
        x = MARGIN_X + col_w * i
        add_text(s, x, y0, col_w - Inches(0.2), Inches(0.9), pct,
                 size=54, bold=True, color=INK)
        add_text(s, x, y0 + Inches(0.95), col_w - Inches(0.2), Inches(0.9), label,
                 size=11, color=INK2, line_spacing=1.4)
        add_text(s, x, y0 + Inches(1.85), col_w - Inches(0.2), Inches(0.3), src.upper(),
                 size=9, bold=True, color=INK3, letter_spacing=160)
    # quote box
    bx_l = MARGIN_X; bx_t = Inches(6.05); bx_w = Inches(11); bx_h = Inches(0.85)
    add_rect(s, bx_l, bx_t, bx_w, bx_h, fill=SURFACE, line=BLUE, line_w=Pt(0.75), radius=Emu(228600))
    add_text(s, bx_l + Inches(0.3), bx_t + Inches(0.15), bx_w - Inches(0.5), bx_h - Inches(0.3),
             "The intent is there. The outcome rarely is. Between the strategy deck and the shop floor — something breaks.",
             size=14, italic=True, color=BLUE, line_spacing=1.3)
    add_footer(s, "Strategy on paper. No execution on the ground.", 3)


def slide_04_unserved(prs):
    s = blank_slide(prs)
    add_chrome(s, "Act 1 · The Unserved Market · 04")
    add_eyebrow(s, Inches(0.85), "The market gap")
    add_text(s, MARGIN_X, Inches(1.5), Inches(11.5), Inches(1.5),
             [[{"t": "Hundreds of thousands of SMEs.", "size": 36, "bold": True}],
              [{"t": "The biggest players still can't fit.", "size": 36, "bold": True, "italic": True, "color": BLUE}]],
             color=INK, line_spacing=1.05)
    add_text(s, MARGIN_X, Inches(3.0), Inches(11.5), Inches(0.8),
             "The large ERP platforms, the Tier-1 consultancies, the global implementation partners — all focus elsewhere. Not because they don't see the SME market. Because their structure cannot serve it profitably.",
             size=13, color=INK2, line_spacing=1.45)

    cols = [
        ("01", "Price floor.",
         "Enterprise consultants charge S$200–550 per hour. Professional services typically run 100–200% of annual software licence fees. A 30-person SME cannot absorb a six-figure bill — and a Tier-1 vendor cannot profitably quote below its own fixed cost base.",
         "Industry rates · Panorama Consulting"),
        ("02", "Product built for someone else.",
         "Modules are built around idealised enterprise workflows: standardised inputs, clean master data, uniform roles. SME operations are messier — varied processes, tribal knowledge, non-standard work. The system does not bend. The SME is told to.",
         "ERP-in-SME literature · MDPI 2021, 2016"),
        ("03", "Incentive misalignment.",
         "Vendor revenue is licence + hosting + annual maintenance — paid whether the client sees value or not. \"Finishing\" is not the goal. Dependency is the business model. A client who finishes and self-manages is a customer lost.",
         "Observed across ERP disasters · CIO"),
    ]
    col_w = CONTENT_W / 3
    y0 = Inches(4.05)
    for i, (idx, h, body, src) in enumerate(cols):
        x = MARGIN_X + col_w * i
        add_text(s, x, y0, col_w - Inches(0.2), Inches(0.6), idx,
                 size=30, bold=True, color=BLUE)
        add_text(s, x, y0 + Inches(0.55), col_w - Inches(0.2), Inches(0.5), h,
                 size=15, bold=True, color=INK)
        add_text(s, x, y0 + Inches(1.0), col_w - Inches(0.2), Inches(1.7), body,
                 size=11, color=INK2, line_spacing=1.45)
        add_text(s, x, y0 + Inches(2.4), col_w - Inches(0.2), Inches(0.3), src.upper(),
                 size=8, bold=True, color=INK3, letter_spacing=160)
    add_footer(s, "Big reach. Wrong structure. Wrong price. Wrong incentive.", 4)


def slide_05_structural(prs):
    s = blank_slide(prs)
    add_chrome(s, "Act 1 · Our Structural Fit · 05")
    add_eyebrow(s, Inches(0.85), "Why we fit — structurally")
    add_text(s, MARGIN_X, Inches(1.5), Inches(11.5), Inches(1.5),
             [[{"t": "Three advantages.", "size": 36, "bold": True}],
              [{"t": "Not strategy. Structure.", "size": 36, "bold": True, "italic": True, "color": BLUE}]],
             color=INK, line_spacing=1.05)
    add_text(s, MARGIN_X, Inches(3.0), Inches(11.5), Inches(0.8),
             "We don't serve the SME market because we decided it was underserved. We serve it because our entire company is built to — and we would be unprofitable trying to do anything else.",
             size=13, color=INK2, line_spacing=1.45)
    cols = [
        ("01", "Lean cost base.",
         "No licence revenue to protect. No office-tower overhead to fund. Specialist engineers sourced globally, where cost-to-capability is strongest. The advantage is passed through to the client — not pocketed as margin. Our S$8,000 fixed scope is not a discount. It is our natural price."),
        ("02", "Engineering fit.",
         "We are engineers. Not IT consultants. Not business graduates with a factory manual. We have wired panels, debugged production lines, walked shop floors at 2am. When we configure an inventory module, we know where the physical bins sit."),
        ("03", "Aligned incentive.",
         "We earn from work delivered. No licence markups, no platform commissions, no annual maintenance contracts. When the work is done, we leave. You own the software, the server, the data, the workflows. That alignment is why SMEs finally trust us to actually finish."),
    ]
    col_w = CONTENT_W / 3
    y0 = Inches(4.05)
    for i, (idx, h, body) in enumerate(cols):
        x = MARGIN_X + col_w * i
        add_text(s, x, y0, col_w - Inches(0.2), Inches(0.6), idx,
                 size=30, bold=True, color=BLUE)
        add_text(s, x, y0 + Inches(0.55), col_w - Inches(0.2), Inches(0.45), h,
                 size=15, bold=True, color=INK)
        add_text(s, x, y0 + Inches(1.0), col_w - Inches(0.2), Inches(1.9), body,
                 size=11, color=INK2, line_spacing=1.45)
    add_footer(s, "Specialist doctor. Not salesman. Built to serve — at a price that works.", 5)


def slide_06_problem1(prs):
    s = blank_slide(prs)
    add_chrome(s, "Act 1 · The Problem · 06")
    add_eyebrow(s, Inches(0.85), "Problem 01 of 03")
    add_text(s, MARGIN_X, Inches(1.5), Inches(11), Inches(1.6),
             [[{"t": "You can't uphold", "size": 38, "bold": True, "color": INK}],
              [{"t": "your own ", "size": 38, "bold": True, "color": INK},
               {"t": "data, knowledge, and history.", "size": 38, "bold": True, "italic": True, "color": BLUE}]],
             line_spacing=1.1)

    rows = [
        ("01 — Can't learn", "The same problem keeps repeating.",
         "Improvement meetings run on bias and loudest voice. Supplier lead times drift for months before anyone notices. Growth stalls at the founder's bandwidth, because every lesson has to pass through the founder's memory."),
        ("02 — People leave", "Heavy penalty every time a key person resigns — or is just unavailable.",
         "Retention cost spirals. You pay above market to keep the one person who knows how things work. A single resignation stalls entire workflows. Project risk is one bad Monday away."),
        ("03 — New hires can't catch up", "Onboarding takes months of tribal knowledge transfer. Most leave before becoming productive.",
         "You end up hiring above market rate just to attract people willing to absorb the chaos — effectively training talent that competitors will later poach."),
    ]
    y = Inches(3.55)
    row_h = Inches(1.1)
    for idx, title, body in rows:
        add_hline(s, MARGIN_X, y, Inches(11))
        add_text(s, MARGIN_X, y + Inches(0.15), Inches(2.2), Inches(0.4), idx.upper(),
                 size=10, bold=True, color=BLUE, letter_spacing=160)
        add_text(s, MARGIN_X + Inches(2.4), y + Inches(0.1), Inches(8.6), Inches(0.45), title,
                 size=15, bold=True, color=INK)
        add_text(s, MARGIN_X + Inches(2.4), y + Inches(0.5), Inches(8.6), Inches(0.6), body,
                 size=11, color=INK2, line_spacing=1.45)
        y += row_h
    add_hline(s, MARGIN_X, y, Inches(11))
    add_footer(s, "Knowledge that lives in one head is knowledge the company doesn't own.", 6)


def slide_07_problem2(prs):
    s = blank_slide(prs)
    add_chrome(s, "Act 1 · The Problem · 07")
    add_eyebrow(s, Inches(0.85), "Problem 02 of 03")
    # Left
    add_text(s, MARGIN_X, Inches(1.55), Inches(5.5), Inches(2.6),
             [[{"t": "You're paying", "size": 34, "bold": True}],
              [{"t": "salary for ", "size": 34, "bold": True},
               {"t": "duplicate work", "size": 34, "bold": True, "italic": True, "color": BLUE}],
              [{"t": "— and you don't track the cost.", "size": 34, "bold": True}]],
             color=INK, line_spacing=1.1)
    add_text(s, MARGIN_X, Inches(4.6), Inches(5.5), Inches(1.5),
             "Staff spend their days transferring information from hardcopy to softcopy, from softcopy to another softcopy, from email to spreadsheet, from spreadsheet to WhatsApp. Zero new information is created.",
             size=13, color=INK2, line_spacing=1.45)

    # Right: cost rows
    rx = MARGIN_X + Inches(6.0)
    rw = Inches(5.0)
    add_text(s, rx, Inches(1.55), rw, Inches(0.3), "THE HIDDEN SALARY BILL",
             size=10, bold=True, color=BLUE, letter_spacing=160)
    rows = [
        ("Purchasing clerk", "Re-types supplier quotations from email PDFs into an Excel tracker.", "22 hrs", "/month"),
        ("Production supervisor", "Copies paper traveller job numbers into a daily WhatsApp report.", "5.5 hrs", "/month"),
        ("Accounts admin", "Copies delivery order numbers from hardcopy binders into accounting software.", "8 hrs", "/month"),
        ("QC inspector", "Logs rejection data on paper, re-enters into quarterly summaries.", "16 hrs", "/qtr"),
    ]
    y = Inches(1.95)
    for role, what, hrs, suf in rows:
        add_hline(s, rx, y, rw)
        add_text(s, rx, y + Inches(0.1), Inches(1.8), Inches(0.35), role,
                 size=12, bold=True, italic=True, color=INK)
        add_text(s, rx + Inches(1.85), y + Inches(0.1), Inches(2.0), Inches(0.65), what,
                 size=10, color=INK2, line_spacing=1.4)
        add_text(s, rx + rw - Inches(1.0), y + Inches(0.05), Inches(1.0), Inches(0.7),
                 [[{"t": hrs, "size": 14, "bold": True, "color": INK}],
                  [{"t": suf, "size": 9, "italic": True, "bold": True, "color": BLUE}]],
                 align=PP_ALIGN.RIGHT, line_spacing=1.1)
        y += Inches(0.78)
    add_hline(s, rx, y, rw)

    # Bottom blue box
    bx_t = Inches(5.6); bx_h = Inches(1.0)
    add_rect(s, rx, bx_t, rw, bx_h, fill=SURFACE2, radius=Emu(228600))
    add_rect(s, rx, bx_t, Pt(3), bx_h, fill=BLUE)
    add_text(s, rx + Inches(0.2), bx_t + Inches(0.18), rw - Inches(0.4), bx_h - Inches(0.3),
             [[
                 {"t": "A 30-person SME typically loses ", "italic": True},
                 {"t": "1–1.5 full-time headcount", "italic": True, "bold": True, "color": BLUE},
                 {"t": " to this work. Every month. Paid for activity that creates no new value.", "italic": True},
             ]],
             size=11, color=INK, line_spacing=1.35)

    add_footer(s, "Delay. Error. Rework. And nobody tracks the bill.", 7)


def slide_08_problem3(prs):
    s = blank_slide(prs)
    add_chrome(s, "Act 1 · The Problem · 08")
    add_eyebrow(s, Inches(0.85), "Problem 03 of 03")
    add_text(s, MARGIN_X, Inches(1.55), Inches(6.0), Inches(2.6),
             [[{"t": "When the data", "size": 34, "bold": True}],
              [{"t": "going in is ", "size": 34, "bold": True},
               {"t": "unreliable", "size": 34, "bold": True, "italic": True, "color": BLUE},
               {"t": "—", "size": 34, "bold": True}],
              [{"t": "everything built on top is unreliable.", "size": 34, "bold": True}]],
             color=INK, line_spacing=1.1)
    add_text(s, MARGIN_X, Inches(4.7), Inches(6.0), Inches(1.6),
             "Inconsistent inputs become inconsistent reports. Inconsistent reports become wrong decisions. Wrong decisions become losses you only discover after the damage is done.",
             size=13, color=INK2, line_spacing=1.45)

    boxes = [
        ("01 · Wasted manpower",
         "Hours per week correcting, reviewing, reconciling — effectively a full-time role cleaning what should have been right at source."),
        ("02 · Wrong decisions",
         "Buying extra stock \"just in case.\" Investing in the wrong equipment. Hiring headcount because workload reports are wrong."),
        ("03 · No deep analysis",
         "Trend analysis, margin analysis, customer profitability — impossible to run on unreliable data."),
        ("04 · AI remains out of reach",
         "AI forecasting, vision QC, predictive maintenance — all require clean structured data. Messy data = no AI = watching competitors compound every year."),
    ]
    rx = MARGIN_X + Inches(6.5); rw = Inches(4.5)
    y = Inches(1.55)
    for i, (h, b) in enumerate(boxes):
        accent = (i == 3)
        box_h = Inches(1.15)
        add_rect(s, rx, y, rw, box_h,
                 fill=SURFACE if accent else SURFACE2,
                 line=BLUE if accent else None,
                 line_w=Pt(0.75) if accent else None,
                 radius=Emu(228600))
        add_text(s, rx + Inches(0.2), y + Inches(0.15), rw - Inches(0.4), Inches(0.3),
                 h.upper(), size=10, bold=True, color=BLUE, letter_spacing=160)
        add_text(s, rx + Inches(0.2), y + Inches(0.5), rw - Inches(0.4), Inches(0.6),
                 b, size=11, color=INK2, line_spacing=1.4)
        y += box_h + Inches(0.1)
    add_footer(s, "The AI era is here. The question is whether your data is ready for it.", 8)


def slide_09_compound(prs):
    s = blank_slide(prs)
    add_chrome(s, "Act 1 · The Compound Damage · 09")
    add_eyebrow(s, Inches(0.85), "What standing still actually costs")
    add_text(s, MARGIN_X, Inches(1.45), Inches(11), Inches(1.4),
             [[{"t": "Not four separate problems.", "size": 32, "bold": True}],
              [{"t": "One chain.", "size": 32, "bold": True, "italic": True, "color": BLUE},
               {"t": " Each link forges the next.", "size": 32, "bold": True}]],
             color=INK, line_spacing=1.1)

    nodes = [
        ("Link 01", "Poor operating efficiency",
         ["Output per employee drops", "Work duplicated and re-done daily", "Friction compounds in every workflow"]),
        ("Link 02", "High burnout & turnover",
         ["Younger staff refuse low-value rework", "Job satisfaction collapses under chaos", "Career anxiety drives good people out"]),
        ("Link 03", "Poor expense efficiency",
         ["Inventory can't be kept lean — cash trapped", "Purchasing reactive, no leverage", "Forecasting breaks down both sides"]),
        ("Link 04", "Losing competitiveness",
         ["Higher cost base loses bids on price", "Slow responses damage customer experience", "Thin margin leaves nothing to reinvest"]),
    ]
    n = 4
    arrow_w = Inches(0.35)
    node_w = (CONTENT_W - 3 * arrow_w) / n
    y = Inches(3.55); h_node = Inches(2.3)
    for i, (lnk, head, items) in enumerate(nodes):
        x = MARGIN_X + (node_w + arrow_w) * i
        add_rect(s, x, y, node_w, h_node, fill=SURFACE, radius=Emu(228600))
        add_rect(s, x, y, node_w, Pt(2), fill=BLUE)
        add_text(s, x + Inches(0.2), y + Inches(0.15), node_w - Inches(0.4), Inches(0.3),
                 lnk.upper(), size=9, bold=True, color=BLUE, letter_spacing=160)
        add_text(s, x + Inches(0.2), y + Inches(0.45), node_w - Inches(0.4), Inches(0.6),
                 head, size=14, bold=True, color=INK, line_spacing=1.15)
        ity = y + Inches(1.1)
        for it in items:
            add_hline(s, x + Inches(0.2), ity, node_w - Inches(0.4))
            add_text(s, x + Inches(0.25), ity + Inches(0.05), node_w - Inches(0.5), Inches(0.35),
                     [[{"t": "— ", "color": BLUE}, {"t": it, "color": INK2}]],
                     size=10, line_spacing=1.4)
            ity += Inches(0.4)
        if i < n - 1:
            add_text(s, x + node_w, y, arrow_w, h_node, "→",
                     size=20, bold=True, italic=True, color=BLUE,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Bottom blue line box
    bx_t = Inches(6.05); bx_h = Inches(0.7)
    add_rect(s, MARGIN_X, bx_t, Inches(11), bx_h, fill=SURFACE2, radius=Emu(228600))
    add_rect(s, MARGIN_X, bx_t, Pt(3), bx_h, fill=BLUE)
    add_text(s, MARGIN_X + Inches(0.25), bx_t + Inches(0.12), Inches(10.5), bx_h - Inches(0.2),
             "The chain runs one way. Break it at Link 01 — operating efficiency — and the rest of the chain never forms.",
             size=12, italic=True, color=BLUE, line_spacing=1.3)

    add_footer(s, "Four cascading damages. One starting point. One lever that breaks the whole chain.", 9)


def slide_10_mindset(prs):
    s = blank_slide(prs)
    add_chrome(s, "Act 1 · The Mindset Tax · 10")
    add_eyebrow(s, Inches(0.85), "The most expensive sentence in the language")
    # Left mega quote
    add_text(s, MARGIN_X, Inches(1.7), Inches(6.0), Inches(5.0),
             [[{"t": "“We've", "size": 78, "bold": True, "italic": True, "color": BLUE}],
              [{"t": "always", "size": 78, "bold": True, "italic": True, "color": BLUE}],
              [{"t": "run it", "size": 78, "bold": True, "italic": True, "color": BLUE}],
              [{"t": "this way.”", "size": 78, "bold": True, "italic": True, "color": BLUE}]],
             line_spacing=0.98)
    # Right column
    rx = MARGIN_X + Inches(6.5); rw = Inches(4.7)
    add_text(s, rx, Inches(1.7), rw, Inches(1.2),
             "Not a warning. An invoice — already accumulating, quietly, every month.",
             size=18, italic=True, bold=True, color=BLUE, line_spacing=1.3)
    add_text(s, rx, Inches(3.0), rw, Inches(2.0),
             "Every SME owner says this at least once. And the longer it's said, the more the cost compounds. The bill doesn't arrive in a letter — it shows up as a lost bid, a resigned foreman, an inventory count that no longer matches reality.",
             size=12, color=INK2, line_spacing=1.45)
    add_text(s, rx, Inches(4.7), rw, Inches(1.4),
             "This is not a question of technology. It is a question of whether the business continues to run on one person's memory — or on a system the company actually owns.",
             size=12, color=INK2, line_spacing=1.45)
    # source box
    add_rect(s, rx, Inches(6.05), rw, Inches(0.7), fill=SURFACE2, radius=Emu(180000))
    add_text(s, rx + Inches(0.2), Inches(6.12), rw - Inches(0.4), Inches(0.3),
             "SOURCES REFERENCED", size=9, bold=True, color=BLUE, letter_spacing=160)
    add_text(s, rx + Inches(0.2), Inches(6.4), rw - Inches(0.4), Inches(0.3),
             "SBF National Business Survey 2025 · EDB / MAS 2026 · SME Development Survey",
             size=10, color=INK2)
    add_footer(s, "Inertia has a monthly price tag — and you're already paying it.", 10)


def slide_11_mentor(prs):
    s = blank_slide(prs)
    add_chrome(s, "Act 2 · Meeting the Mentor · 11")
    add_eyebrow(s, Inches(0.85), "Who we are")
    # Left headline
    add_text(s, MARGIN_X, Inches(1.55), Inches(6.0), Inches(5.0),
             [[{"t": "We don't sell", "size": 38, "bold": True, "color": INK}],
              [{"t": "software.", "size": 38, "bold": True, "color": INK}],
              [{"t": "We help SMEs", "size": 38, "bold": True, "italic": True, "color": BLUE}],
              [{"t": "own the system", "size": 38, "bold": True, "italic": True, "color": BLUE}],
              [{"t": "that runs their", "size": 38, "bold": True, "italic": True, "color": BLUE}],
              [{"t": "business.", "size": 38, "bold": True, "italic": True, "color": BLUE}]],
             line_spacing=1.05)
    # Right three points
    rx = MARGIN_X + Inches(6.5); rw = Inches(4.7)
    points = [
        ("01", "Engineers only.",
         "Not IT consultants reading a factory manual. Engineers who have wired control panels, pulled wire harnesses, walked production floors, debugged lines at 2am."),
        ("02", "Zero conflict of interest.",
         "We don't earn from software licences. We don't mark up hosting. We don't take platform commissions. We earn from work delivered. Nothing else."),
        ("03", "You own everything.",
         "Your software. Your server. Your data. Your workflows. When our work is done, you don't need us. That is the point."),
    ]
    y = Inches(1.55)
    for idx, head, body in points:
        ring = slide_ring(s, rx, y, idx)
        add_text(s, rx + Inches(0.7), y, rw - Inches(0.7), Inches(0.4), head,
                 size=18, bold=True, color=INK)
        add_text(s, rx + Inches(0.7), y + Inches(0.45), rw - Inches(0.7), Inches(1.1), body,
                 size=11, color=INK2, line_spacing=1.45)
        y += Inches(1.65)
    add_footer(s, "Specialist doctor. Not salesman.", 11)


def slide_ring(slide, x, y, label):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.5), Inches(0.5))
    shp.fill.background()
    shp.line.color.rgb = BLUE
    shp.line.width = Pt(1.25)
    tf = shp.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.name = FONT; r.font.size = Pt(11); r.font.bold = True
    r.font.color.rgb = BLUE
    return shp


def slide_12_question(prs):
    s = blank_slide(prs)
    add_chrome(s, "Act 2 · What They Really Ask · 12")
    add_eyebrow(s, Inches(0.85), "A moment every SME owner has with us")
    add_text(s, MARGIN_X, Inches(1.5), Inches(11), Inches(1.5),
             "Every SME owner asks us the same question. Eventually. Usually not in the first meeting — and not on a slide. It comes later, when the deck is closed and we're being walked to the door. Voice lowered. Almost casual. Like it's a secret they've been carrying for two years.",
             size=15, italic=True, color=INK2, line_spacing=1.45)
    add_text(s, MARGIN_X, Inches(3.3), Inches(11), Inches(2.5),
             [[{"t": "“What happens", "size": 64, "bold": True, "italic": True, "color": BLUE}],
              [{"t": "when you leave?”", "size": 64, "bold": True, "italic": True, "color": BLUE}]],
             line_spacing=1.05)
    add_hline(s, MARGIN_X, Inches(6.0), Inches(11))
    add_text(s, MARGIN_X, Inches(6.15), Inches(2.4), Inches(0.3),
             "WHAT THEY'RE REALLY ASKING", size=9, bold=True, color=INK3, letter_spacing=160)
    add_text(s, MARGIN_X + Inches(2.6), Inches(6.05), Inches(8.4), Inches(0.6),
             "\"Am I going to end up stuck — the way I did with the last vendor?\"",
             size=18, italic=True, bold=True, color=INK)
    add_footer(s, "It's the only question that matters. The whole offer lives or dies on the answer.", 12)


def slide_13_six_answers(prs):
    s = blank_slide(prs)
    add_chrome(s, "Act 2 · The Six Answers · 13")
    add_eyebrow(s, Inches(0.85), "Six questions. Six structural answers.")
    add_text(s, MARGIN_X, Inches(1.45), Inches(11), Inches(0.9),
             [[{"t": "Not a piece of software. ", "size": 22, "bold": True, "color": INK},
               {"t": "A way of operating", "size": 22, "bold": True, "italic": True, "color": BLUE},
               {"t": " — designed around the questions every SME owner actually asks.", "size": 22, "bold": True, "color": INK}]],
             line_spacing=1.2)
    qa = [
        ("01", "\"What happens to my data — if we ever part ways?\"",
         "One integrated database, on your server. When we leave, you own every byte — the records, the history, the lessons."),
        ("02", "\"What if my staff forgets how to use it?\"",
         "The workflow is the system. No training PDF to forget, no memo to lose — the right way of working becomes the default way."),
        ("03", "\"What if I need to change something later?\"",
         "Drag-and-drop configuration. Your own staff extends and adjusts it. No support tickets. No annual contract."),
        ("04", "\"What if cash flow gets tight for a few months?\"",
         "No licence fees. No monthly subscription trap. You pay for tasks completed — and you can pause anytime."),
        ("05", "\"What if we outgrow it?\"",
         "Modular. Plug in one piece at a time. Keep what already works. Never rip-and-replace."),
        ("06", "\"How does the team actually stay coordinated through all this?\"",
         "One portal. Every message, every approval, every change tied to a record. One place to look — no more searching five apps for one answer."),
    ]
    y = Inches(2.65)
    row_h = Inches(0.65)
    add_hline(s, MARGIN_X, y, Inches(11))
    for num, q, a in qa:
        add_text(s, MARGIN_X, y + Inches(0.12), Inches(0.5), Inches(0.4), num,
                 size=10, bold=True, color=BLUE, letter_spacing=160)
        add_text(s, MARGIN_X + Inches(0.55), y + Inches(0.1), Inches(5.0), Inches(0.55), q,
                 size=12, italic=True, color=INK3, line_spacing=1.3)
        add_hline(s, MARGIN_X + Inches(5.7), y, Pt(0.5))  # subtle separator
        # vertical separator
        sep = s.shapes.add_connector(1, MARGIN_X + Inches(5.7), y + Inches(0.08),
                                     MARGIN_X + Inches(5.7), y + row_h - Inches(0.05))
        sep.line.color.rgb = RULE; sep.line.width = Pt(0.5)
        add_text(s, MARGIN_X + Inches(5.85), y + Inches(0.1), Inches(5.15), Inches(0.55), a,
                 size=12, color=INK, line_spacing=1.3)
        y += row_h
        add_hline(s, MARGIN_X, y, Inches(11))
    add_footer(s, "Every answer is a structural commitment — not a promise.", 13)


def slide_14_services(prs):
    s = blank_slide(prs)
    add_chrome(s, "Act 2 · The Service · 14")
    add_eyebrow(s, Inches(0.85), "What we deliver")
    add_text(s, MARGIN_X, Inches(1.5), Inches(11), Inches(1.4),
             [[{"t": "Three layers. ", "size": 36, "bold": True, "color": INK},
               {"t": "Fully defined scope.", "size": 36, "bold": True, "italic": True, "color": BLUE}],
              [{"t": "One accountable partner.", "size": 36, "bold": True, "color": INK}]],
             line_spacing=1.1)
    cols = [
        ("Layer 01 · Software", "Seven core module areas.",
         ["Procurement & Purchasing", "Inventory & Warehouse", "Manufacturing & Production",
          "Quality Management", "Finance & Accounting", "HR & Payroll", "Sales & CRM"],
         "Additional modules (Safety, Logistics, Projects) plugged in as needed."),
        ("Layer 02 · Implementation Labour", "The human work that makes software actually work.",
         ["Implementation, configuration, setup", "Role-specific on-site training",
          "Data migration (clean, map, validate)", "Hardcopy-to-softcopy data entry",
          "On-site survey and data capture", "Scheduled reporting setup", "Per-task ongoing fixes"],
         "On-demand human work — scoped task by task."),
        ("Layer 03 · Customisation", "For clients whose needs exceed standard modules.",
         ["Integration with existing software", "R&D of new software modules",
          "Hardware sourcing", "Control panel design & build",
          "Wire harness & PCBA assembly", "IoT sensor deployment", "Custom fixtures"],
         "Engineering DNA — not outsourced, not guessed."),
    ]
    col_w = (CONTENT_W - Inches(0.4)) / 3
    y = Inches(3.65); h = Inches(3.3)
    for i, (tag, head, items, foot) in enumerate(cols):
        x = MARGIN_X + (col_w + Inches(0.2)) * i
        add_rect(s, x, y, col_w, h, fill=SURFACE, radius=Emu(228600))
        ix = x + Inches(0.3); iw = col_w - Inches(0.6)
        add_text(s, ix, y + Inches(0.25), iw, Inches(0.3), tag.upper(),
                 size=10, bold=True, color=BLUE, letter_spacing=160)
        add_text(s, ix, y + Inches(0.6), iw, Inches(0.7), head,
                 size=14, bold=True, color=INK, line_spacing=1.15)
        iy = y + Inches(1.35)
        for it in items:
            add_text(s, ix, iy, iw, Inches(0.28), it, size=10, color=INK2)
            iy += Inches(0.26)
            add_hline(s, ix, iy, iw, color=RULE, weight=Pt(0.5))
        add_text(s, ix, y + h - Inches(0.5), iw, Inches(0.4), foot,
                 size=9, color=INK3, italic=True, line_spacing=1.4)
    add_footer(s, "Software. Labour. Customisation. Three layers. One accountable partner.", 14)


def slide_15_fabrication(prs):
    s = blank_slide(prs)
    add_chrome(s, "Act 2 · Fabrication & Production · 15")
    add_eyebrow(s, Inches(0.85), "When the system needs hardware, not just software")
    add_text(s, MARGIN_X, Inches(1.5), Inches(11), Inches(1.4),
             [[{"t": "Production hardware.", "size": 32, "bold": True, "color": INK}],
              [{"t": "Sourced, managed, delivered.", "size": 32, "bold": True, "italic": True, "color": BLUE}]],
             line_spacing=1.1)
    add_text(s, MARGIN_X, Inches(2.95), Inches(11), Inches(0.8),
             "One partner. A curated vendor network across Singapore, China, and India. We scope the requirement, select the right vendor, manage the schedule, run quality assurance, and arrange delivery — so you deal with us, not three factories.",
             size=12, color=INK2, line_spacing=1.45)
    cats = [
        ("A · Metal Fabrication", "Custom structures, enclosures, assemblies.",
         ["Steel structures & assemblies", "Electrical enclosures", "Architectural aluminium & cladding",
          "Rail transit components", "EV / ESU housings"]),
        ("B · Cable & Harness", "Single harness to batch production.",
         ["Industrial wire harnesses", "Automotive-grade assemblies", "Defense-grade cabling",
          "Custom power cables", "Fiber optic & defense optics"]),
        ("C · Control Panels", "Relay logic to PLC automation.",
         ["Industrial & remote I/O panels", "Stainless panels for harsh sites",
          "Pre-wired box-build assemblies", "Ex-rated junction boxes", "Schneider PLC · Magelis HMI"]),
        ("D · PCB & Electronics", "Board assembly & box-build.",
         ["SMT PCB assembly", "Through-hole & mixed-tech", "Box-build electronics",
          "Industrial & defense-grade", "Medical, biometric, railway"]),
        ("E · New Energy", "EV, ESU, and grid-tied hardware.",
         ["EV charging cabinets", "ESU control cabinets", "Lithium battery oven housings",
          "Grid-tied structural components", "Custom enclosures, to spec"]),
    ]
    col_w = (CONTENT_W - Inches(0.4)) / 5
    y = Inches(4.0); h = Inches(2.4)
    for i, (tag, head, items) in enumerate(cats):
        x = MARGIN_X + (col_w + Inches(0.1)) * i
        add_rect(s, x, y, col_w, h, fill=SURFACE, radius=Emu(180000))
        ix = x + Inches(0.15); iw = col_w - Inches(0.3)
        add_text(s, ix, y + Inches(0.15), iw, Inches(0.3), tag.upper(),
                 size=8, bold=True, color=BLUE, letter_spacing=160)
        add_text(s, ix, y + Inches(0.45), iw, Inches(0.7), head,
                 size=11, bold=True, color=INK, line_spacing=1.15)
        iy = y + Inches(1.15)
        for it in items:
            add_text(s, ix, iy, iw, Inches(0.22), it, size=8, color=INK2, line_spacing=1.3)
            iy += Inches(0.22)

    # Bottom dashed accent line — render as outlined rect
    bx_t = Inches(6.55); bx_h = Inches(0.55)
    bx = add_rect(s, MARGIN_X, bx_t, Inches(11), bx_h, fill=None, line=BLUE, line_w=Pt(0.75))
    add_text(s, MARGIN_X + Inches(0.3), bx_t + Inches(0.1), Inches(8), bx_h - Inches(0.2),
             [[{"t": "Singapore quality, global cost.", "italic": True, "color": BLUE, "bold": True},
               {"t": " One scope. One schedule. One QA standard. One point of accountability — local.",
                "color": INK2}]],
             size=11, line_spacing=1.3)
    add_text(s, MARGIN_X + Inches(8.5), bx_t + Inches(0.15), Inches(2.5), bx_h - Inches(0.2),
             "FIVE CATEGORIES · ONE PARTNER",
             size=9, bold=True, color=BLUE, letter_spacing=160, align=PP_ALIGN.RIGHT)
    add_footer(s, "Not a marketplace. A vetted vendor network — managed by engineers.", 15)


def slide_16_positioning(prs):
    s = blank_slide(prs)
    add_chrome(s, "Act 2 · The Positioning · 16")
    add_eyebrow(s, Inches(0.85), "Why us, instead of everyone else")
    add_text(s, MARGIN_X, Inches(1.5), Inches(11), Inches(1.4),
             [[{"t": "Three structural differences", "size": 34, "bold": True, "color": INK}],
              [{"t": "incumbents ", "size": 34, "bold": True, "color": INK},
               {"t": "cannot copy.", "size": 34, "bold": True, "italic": True, "color": BLUE}]],
             line_spacing=1.1)
    cols = [
        ("01", "Honest pricing, by structure.",
         "Others earn from licence fees, annual maintenance, and expensive IT day rates — recurring income whether or not you see value. We earn from work delivered. Our monthly fee is deliberately low, or zero between projects. Finishing the work is the objective."),
        ("02", "Lean company, global talent.",
         "No complicated internal workflow. No heavy overhead. We source specialist engineers worldwide, passing the cost advantage through to you — not pocketing it. Large consultancies can't quote S$8K for a fixed scope because their structure makes it unprofitable. Ours doesn't."),
        ("03", "Engineers only.",
         "Not IT consultants with a certification. Not business graduates reading a factory manual. Engineers who have wired panels, written PLC code, debugged production lines. When we configure your inventory module, we know where the physical bins sit."),
    ]
    col_w = CONTENT_W / 3
    y = Inches(3.5)
    for i, (idx, h, body) in enumerate(cols):
        x = MARGIN_X + col_w * i
        add_text(s, x, y, col_w - Inches(0.3), Inches(0.7), idx,
                 size=30, bold=True, color=BLUE)
        add_text(s, x, y + Inches(0.65), col_w - Inches(0.3), Inches(0.5), h,
                 size=15, bold=True, color=INK)
        add_text(s, x, y + Inches(1.15), col_w - Inches(0.3), Inches(2.0), body,
                 size=11, color=INK2, line_spacing=1.45)

    # bottom blue line box
    bx_t = Inches(6.05); bx_h = Inches(0.85)
    add_rect(s, MARGIN_X, bx_t, Inches(11), bx_h, fill=SURFACE2, radius=Emu(180000))
    add_rect(s, MARGIN_X, bx_t, Pt(3), bx_h, fill=BLUE)
    add_text(s, MARGIN_X + Inches(0.25), bx_t + Inches(0.12), Inches(10.5), bx_h - Inches(0.2),
             [[{"t": "The question isn't ", "italic": True, "color": INK},
               {"t": "\"how can a non-IT company configure systems?\" ", "italic": True, "color": BLUE},
               {"t": "The question is — ", "italic": True, "color": INK},
               {"t": "why would you want someone who's never touched a machine configuring the system that runs your factory?",
                "italic": True, "color": BLUE}]],
             size=11, line_spacing=1.35)
    add_footer(s, "Strike where they are weak, not where they are strong.", 16)


def slide_17_before(prs):
    s = blank_slide(prs)
    add_chrome(s, "Act 3 · Case Study · Before · 17")
    add_eyebrow(s, Inches(0.85), "The prospect's starting state")
    # Left
    add_text(s, MARGIN_X, Inches(1.55), Inches(5.5), Inches(2.8),
             [[{"t": "A 40-person", "size": 34, "bold": True, "color": INK}],
              [{"t": "precision machining", "size": 34, "bold": True, "color": INK}],
              [{"t": "shop in Tuas.", "size": 34, "bold": True, "color": INK}],
              [{"t": "Twelve years running.", "size": 34, "bold": True, "italic": True, "color": BLUE}]],
             line_spacing=1.05)
    add_text(s, MARGIN_X, Inches(4.7), Inches(5.5), Inches(0.95),
             "Spent S$65K on an ERP two years ago. Vendor ran impressive demos, promised the world, disappeared six months after go-live.",
             size=12, color=INK2, line_spacing=1.45)
    add_text(s, MARGIN_X, Inches(5.7), Inches(5.5), Inches(0.95),
             "Today only 2 of 12 modules are used. Staff bypass the system for Excel. Inventory data doesn't match the warehouse.",
             size=12, color=INK2, line_spacing=1.45)

    # Right stats
    rx = MARGIN_X + Inches(6.0); rw = Inches(5.0)
    stats = [
        ("Module adoption", "2 of 12 active", "Staff work around the system, not with it.", False),
        ("Inventory accuracy", "~60%", "System says one thing. The shelf says another.", False),
        ("Monthly vendor cost", "S$2,500/mo", "Hosting + \"support.\" Tickets take 5 days.", False),
    ]
    y = Inches(1.55)
    for k, v, note, accent in stats:
        h_box = Inches(0.85)
        add_rect(s, rx, y, rw, h_box, fill=SURFACE, radius=Emu(180000))
        add_text(s, rx + Inches(0.25), y + Inches(0.12), Inches(2.0), Inches(0.3), k.upper(),
                 size=9, bold=True, color=INK3, letter_spacing=160)
        add_text(s, rx + Inches(0.25), y + Inches(0.4), Inches(2.5), Inches(0.4), v,
                 size=18, bold=True, color=INK)
        add_text(s, rx + Inches(2.8), y + Inches(0.18), rw - Inches(3.0), h_box - Inches(0.3), note,
                 size=10, color=INK3, line_spacing=1.4, align=PP_ALIGN.RIGHT)
        y += h_box + Inches(0.1)
    # accent
    h_box = Inches(1.15)
    add_rect(s, rx, y, rw, h_box, fill=SURFACE, line=BLUE, line_w=Pt(0.75), radius=Emu(180000))
    add_text(s, rx + Inches(0.25), y + Inches(0.15), rw - Inches(0.5), Inches(0.3),
             "THE REAL DAMAGE", size=9, bold=True, color=BLUE, letter_spacing=160)
    add_text(s, rx + Inches(0.25), y + Inches(0.45), rw - Inches(0.5), h_box - Inches(0.55),
             "The boss has lost trust in the idea that technology can work for a company his size.",
             size=12, italic=True, color=BLUE, line_spacing=1.3)
    add_footer(s, "This is not a story about software. It's a story about being abandoned mid-implementation.", 17)


def slide_18_after(prs):
    s = blank_slide(prs)
    add_chrome(s, "Act 3 · Case Study · After · 18")
    add_eyebrow(s, Inches(0.85), "90 days later")
    add_text(s, MARGIN_X, Inches(1.5), Inches(11), Inches(2.0),
             [[{"t": "Same software.", "size": 34, "bold": True, "color": INK}],
              [{"t": "Same team.", "size": 34, "bold": True, "color": INK}],
              [{"t": "Completely different outcome.", "size": 34, "bold": True, "italic": True, "color": BLUE}]],
             line_spacing=1.05)

    rows = [
        ("Module adoption", "2 of 12 active", "8 of 12 actively used"),
        ("Inventory accuracy", "~60%", "95%+"),
        ("Monthly system cost", "S$2,500 / month", "S$25 / month"),
        ("Report generation", "Manual, 2 hrs/week", "Automated, real-time"),
        ("Annual hosting savings", "—", "~S$29,700 / year"),
        ("Total rescue cost", "—", "S$12K one-time"),
    ]
    y = Inches(4.05)
    add_hline(s, MARGIN_X, y, Inches(11))
    for k, b, a in rows:
        row_h = Inches(0.42)
        add_text(s, MARGIN_X + Inches(0.1), y + Inches(0.08), Inches(2.5), row_h, k.upper(),
                 size=9, bold=True, color=INK3, letter_spacing=160)
        add_text(s, MARGIN_X + Inches(3.0), y + Inches(0.05), Inches(3.5), row_h, b,
                 size=12, italic=True, color=INK3)
        add_text(s, MARGIN_X + Inches(6.8), y + Inches(0.05), Inches(4.2), row_h, a,
                 size=13, bold=True, color=BLUE)
        y += row_h
        add_hline(s, MARGIN_X, y, Inches(11))
    add_footer(s, "He opens one screen Monday morning. Inventory. Production. Deliveries. Everything.", 18)


def slide_19_goal(prs):
    s = blank_slide(prs)
    add_chrome(s, "Act 3 · The Ultimate Goal · 19")
    add_eyebrow(s, Inches(0.85), "The real prize")
    add_text(s, MARGIN_X, Inches(1.5), Inches(11), Inches(1.7),
             [[{"t": "Every task assigned", "size": 30, "bold": True, "color": INK}],
              [{"t": "to the ", "size": 30, "bold": True, "color": INK},
               {"t": "right mechanism.", "size": 30, "bold": True, "italic": True, "color": BLUE}],
              [{"t": "Nobody wasted on the wrong work.", "size": 30, "bold": True, "color": INK}]],
             line_spacing=1.1)

    cols = [
        ("Mechanism 01", "Human",
         "Judgement, relationships, creative problem-solving, handling exceptions and edge cases.",
         "Repetitive data entry, rule-based routing, status tracking, copying numbers between screens."),
        ("Mechanism 02", "Digital System",
         "Consistent rule-based workflows, record-keeping, scheduled reporting, cross-department visibility.",
         "Interpreting ambiguous situations, moments requiring human judgement and accountability."),
        ("Mechanism 03", "AI",
         "Pattern recognition, forecasting, classification, natural language — but only once clean data exists.",
         "Anything requiring accountability without oversight, or situations with no training data to learn from."),
    ]
    col_w = (CONTENT_W - Inches(0.4)) / 3
    y = Inches(3.7); h = Inches(2.9)
    for i, (tag, head, best, worst) in enumerate(cols):
        x = MARGIN_X + (col_w + Inches(0.2)) * i
        add_rect(s, x, y, col_w, h, fill=SURFACE, radius=Emu(228600))
        ix = x + Inches(0.3); iw = col_w - Inches(0.6)
        add_text(s, ix, y + Inches(0.25), iw, Inches(0.3), tag.upper(),
                 size=10, bold=True, color=BLUE, letter_spacing=160)
        add_text(s, ix, y + Inches(0.55), iw, Inches(0.5), head,
                 size=22, bold=True, color=INK)
        add_text(s, ix, y + Inches(1.15), iw, Inches(0.3), "BEST FOR",
                 size=9, bold=True, color=INK3, letter_spacing=160)
        add_text(s, ix, y + Inches(1.4), iw, Inches(0.6), best,
                 size=10, color=INK2, line_spacing=1.45)
        add_text(s, ix, y + Inches(2.05), iw, Inches(0.3), "WORST FOR",
                 size=9, bold=True, color=INK3, letter_spacing=160)
        add_text(s, ix, y + Inches(2.3), iw, Inches(0.6), worst,
                 size=10, color=INK2, line_spacing=1.45)
    add_footer(s, "Sustainable business. Easier life. AI-ready in 6–12 months — naturally.", 19)


def slide_20_cta(prs):
    s = blank_slide(prs)
    add_chrome(s, "Act 3 · One Next Step · 20")
    add_eyebrow(s, Inches(0.85), "One small step. Zero commitment.")
    add_text(s, MARGIN_X, Inches(1.5), Inches(6.0), Inches(4.5),
             [[{"t": "Book a", "size": 60, "bold": True, "color": INK}],
              [{"t": "30-minute", "size": 60, "bold": True, "color": INK}],
              [{"t": "diagnostic.", "size": 60, "bold": True, "italic": True, "color": BLUE}]],
             line_spacing=1.05)
    rx = MARGIN_X + Inches(6.5); rw = Inches(4.8)
    add_text(s, rx, Inches(1.5), rw, Inches(1.4),
             "Not a sales call. We walk your floor (or your screens), quantify what the current operation is actually costing you — in hours, dollars, risk — and tell you straight whether a system is worth it.",
             size=14, color=INK2, line_spacing=1.45)
    add_text(s, rx, Inches(3.3), rw, Inches(0.7),
             "If it's not worth it for your business — we'll say so.",
             size=14, italic=True, bold=True, color=BLUE, line_spacing=1.4)

    # CTA pill
    btn_w = Inches(2.8); btn_h = Inches(0.55)
    btn = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, rx, Inches(4.1), btn_w, btn_h)
    btn.adjustments[0] = 0.5
    btn.fill.solid(); btn.fill.fore_color.rgb = BLUE
    btn.line.fill.background()
    tf = btn.text_frame
    tf.margin_left = Inches(0.15); tf.margin_right = Inches(0.15)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "Book the diagnostic  →"
    r.font.name = FONT; r.font.size = Pt(14); r.font.bold = False
    r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    add_hline(s, rx, Inches(5.0), rw)
    # contact
    contacts = [
        ("FOUNDER", "Ronnie Yap"),
        ("WHATSAPP / PHONE", "+65 9894 3004"),
        ("EMAIL", "ronnie.yap@veloworking.com"),
    ]
    col_w = rw / 3
    for i, (k, v) in enumerate(contacts):
        x = rx + col_w * i
        add_text(s, x, Inches(5.15), col_w - Inches(0.1), Inches(0.3), k,
                 size=8, bold=True, color=INK3, letter_spacing=160)
        size = 13 if i < 2 else 10
        add_text(s, x, Inches(5.5), col_w - Inches(0.1), Inches(0.5), v,
                 size=size, bold=True, color=INK)

    add_footer(s, "Built for the ones still standing.", 20)


def build():
    prs = make_presentation()
    slide_00_cover(prs)
    slide_01_about(prs)
    slide_02_economy(prs)
    slide_03_paradox(prs)
    slide_04_unserved(prs)
    slide_05_structural(prs)
    slide_06_problem1(prs)
    slide_07_problem2(prs)
    slide_08_problem3(prs)
    slide_09_compound(prs)
    slide_10_mindset(prs)
    slide_11_mentor(prs)
    slide_12_question(prs)
    slide_13_six_answers(prs)
    slide_14_services(prs)
    slide_15_fabrication(prs)
    slide_16_positioning(prs)
    slide_17_before(prs)
    slide_18_after(prs)
    slide_19_goal(prs)
    slide_20_cta(prs)
    out = "/home/user/veloworkingwebsite/asset/Velo-Working-Pitch-Deck-v3.pptx"
    prs.save(out)
    print(f"Saved: {out} ({len(prs.slides.__iter__.__self__._sldIdLst) if False else ''})")


if __name__ == "__main__":
    build()
